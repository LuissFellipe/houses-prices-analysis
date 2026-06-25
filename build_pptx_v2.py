"""
Apresentação v2 — California Housing Storytelling
Design: California Sunset — âmbar, coral, oceano profundo
Assinatura: barra de proximidade oceano em cada slide de zona
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd, numpy as np, io
from pathlib import Path
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.ticker as mticker
from matplotlib.patches import FancyBboxPatch
from matplotlib.colors import LinearSegmentedColormap

# ══════════════════════════════════════════════════════════════════════
# TOKENS DE DESIGN — California Sunset
# ══════════════════════════════════════════════════════════════════════
BG       = RGBColor(0x0E, 0x11, 0x17)   # quase-preto azulado
BG2      = RGBColor(0x18, 0x1F, 0x2E)   # card escuro
BG3      = RGBColor(0x1E, 0x28, 0x3C)   # card médio
GOLD     = RGBColor(0xF5, 0xA6, 0x23)   # âmbar califórnia — destaque principal
OCEAN    = RGBColor(0x0E, 0xA5, 0xE9)   # azul oceano pacífico
CORAL    = RGBColor(0xF9, 0x73, 0x16)   # coral pôr do sol
SAGE     = RGBColor(0x34, 0xD3, 0x99)   # verde-sage — positivo
RED      = RGBColor(0xF4, 0x3F, 0x5E)   # vermelho — alerta
MIST     = RGBColor(0x94, 0xA3, 0xB8)   # cinza névoa — corpo
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
SAND     = RGBColor(0xFE, 0xF3, 0xC7)   # areia dourada — highlight

# Cores por zona (temáticas)
ZONE_COLOR = {
    'INLAND':     (0xF5, 0xA6, 0x23),   # âmbar — seco, deserto, sol forte
    '<1H OCEAN':  (0x34, 0xD3, 0x99),   # sage — equilíbrio, subúrbio verde
    'NEAR OCEAN': (0x0E, 0xA5, 0xE9),   # oceano — litoral
    'NEAR BAY':   (0xA7, 0x8B, 0xFA),   # violeta — urbano premium
    'ISLAND':     (0xF4, 0x3F, 0x5E),   # coral vermelho — exclusivo
}

def rgb(*args):
    if len(args) == 3: return RGBColor(*args)
    return RGBColor(*args[0])

def zc(zona): return rgb(*ZONE_COLOR[zona])

W, H = Inches(13.33), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]

plt.rcParams.update({'font.family':'DejaVu Sans','figure.facecolor':'none'})

# ══════════════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════════════
def slide():
    s = prs.slides.add_slide(BLANK)
    f = s.background.fill; f.solid(); f.fore_color.rgb = BG
    return s

def rect(s, l, t, w, h, color, alpha_ignored=None):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.line.fill.background()
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    return sh

def txt(s, text, l, t, w, h, size=11, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def fig2slide(s, fig, l, t, w, h):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', bbox_inches='tight', dpi=160,
                facecolor='none', transparent=True)
    buf.seek(0)
    s.shapes.add_picture(buf, Inches(l), Inches(t), Inches(w), Inches(h))
    plt.close(fig)

def bar_left(s, t, color=GOLD, h=6.8):
    """Barra vertical esquerda assinatura."""
    rect(s, 0, t, 0.07, h, color)

def eyebrow(s, text, l=0.3, t=0.22, color=GOLD):
    """Label pequeno acima do título."""
    txt(s, text.upper(), l, t, 8, 0.3, size=9, bold=True, color=color)

def title(s, t1, t2=None, l=0.3, t=0.5, size1=30, size2=None, color1=WHITE, color2=GOLD):
    txt(s, t1, l, t, 12.5, 0.65, size=size1, bold=True, color=color1)
    if t2:
        txt(s, t2, l, t + size1*0.018, 12.5, 0.55,
            size=size2 or size1, bold=True, color=color2)

def divider(s, t, l=0.3, w=12.7, color=BG3):
    rect(s, l, t, w, 0.035, color)

def tag(s, text, l, t, bg=GOLD, fg=BG, size=9):
    w = max(len(text) * 0.092 + 0.25, 1.0)
    rect(s, l, t, w, 0.27, bg)
    txt(s, text, l + 0.07, t + 0.03, w - 0.1, 0.24,
        size=size, bold=True, color=fg, align=PP_ALIGN.CENTER)

def kpi(s, l, t, w, h, value, label, color=GOLD, sub=None, size_val=26):
    rect(s, l, t, w, h, BG2)
    rect(s, l, t, 0.055, h, color)
    txt(s, value, l+0.14, t+0.08, w-0.2, 0.55, size=size_val, bold=True, color=color)
    txt(s, label, l+0.14, t+0.62, w-0.2, 0.32, size=9, color=MIST)
    if sub:
        txt(s, sub, l+0.14, t+0.9, w-0.2, 0.28, size=8, color=MIST, italic=True)

def callout(s, l, t, w, h, body, title_txt=None, color=GOLD):
    rect(s, l, t, w, h, BG2)
    rect(s, l, t, 0.055, h, color)
    y = t + 0.1
    if title_txt:
        txt(s, title_txt, l+0.14, y, w-0.22, 0.32, size=10, bold=True, color=color)
        y += 0.38
    txt(s, body, l+0.14, y, w-0.22, h - (y-t) - 0.1, size=9.5, color=MIST)

def zone_badge(s, zona, l=0.3, t=6.9):
    """Barra de proximidade ao oceano — assinatura visual recorrente."""
    zonas = ['INLAND', '<1H OCEAN', 'NEAR OCEAN', 'NEAR BAY', 'ISLAND']
    labels = ['Interior', '<1h Oceano', 'Próx. Oceano', 'Próx. Baía', 'Ilha']
    bw = 12.7 / len(zonas)
    for i, (z, lb) in enumerate(zip(zonas, labels)):
        c = rgb(*ZONE_COLOR[z])
        if z == zona:
            rect(s, 0.3 + i*bw, t, bw, 0.38, c)
            txt(s, lb, 0.3 + i*bw + 0.05, t+0.06, bw-0.08, 0.28,
                size=8, bold=True, color=BG, align=PP_ALIGN.CENTER)
        else:
            rect(s, 0.3 + i*bw, t, bw, 0.38, BG3)
            txt(s, lb, 0.3 + i*bw + 0.05, t+0.06, bw-0.08, 0.28,
                size=7, color=MIST, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════
# DADOS
# ══════════════════════════════════════════════════════════════════════
df = pd.read_csv('housing.csv')
df.columns = ['lon','lat','idade','total_q','total_d','pop','dom','renda','valor','zona']
df['total_d'] = df['total_d'].fillna(df['total_d'].median())
df['qpd']  = df['total_q'] / df['dom']
df['ppd']  = df['pop'] / df['dom']
df['score'] = df['renda'] / (df['valor'] / 100000)

ZS = df.groupby('zona').agg(
    n=('valor','count'), vm=('valor','median'), vmean=('valor','mean'),
    vp25=('valor',lambda x: x.quantile(.25)), vp75=('valor',lambda x: x.quantile(.75)),
    rm=('renda','median'), im=('idade','median'),
    qm=('qpd','median'), pm=('ppd','median'), sc=('score','median'),
).sort_values('vm')

ORDER = ['INLAND','<1H OCEAN','NEAR OCEAN','NEAR BAY','ISLAND']

# ══════════════════════════════════════════════════════════════════════
# FUNÇÕES DE GRÁFICO
# ══════════════════════════════════════════════════════════════════════
BG_HEX   = '#0E1117'
BG2_HEX  = '#181F2E'
MIST_HEX = '#94A3B8'
ZONE_HEX = {k: '#{:02X}{:02X}{:02X}'.format(*v) for k,v in ZONE_COLOR.items()}

def styled_ax(ax, title=None):
    ax.set_facecolor('none')
    for sp in ax.spines.values(): sp.set_visible(False)
    ax.tick_params(colors=MIST_HEX, labelsize=8)
    ax.yaxis.grid(True, color='#1E283C', lw=0.6); ax.set_axisbelow(True)
    if title: ax.set_title(title, color='white', fontsize=10, fontweight='bold', pad=6)

def make_map(col, cmap, title, figsize=(6.5,5.2)):
    fig, ax = plt.subplots(figsize=figsize, facecolor='none')
    sc = ax.scatter(df['lon'], df['lat'], c=df[col], cmap=cmap,
                    s=df['pop']/500, alpha=0.45, linewidths=0)
    cb = fig.colorbar(sc, ax=ax, shrink=0.7, pad=0.02)
    cb.ax.tick_params(colors=MIST_HEX, labelsize=7)
    ax.set_facecolor('none')
    ax.tick_params(colors=MIST_HEX, labelsize=7)
    for sp in ax.spines.values(): sp.set_color('#1E283C')
    ax.set_title(title, color='white', fontsize=10, fontweight='bold', pad=5)
    plt.tight_layout(pad=0.3)
    return fig

def make_zone_scatter(zona_focus, figsize=(6.0,4.5)):
    """Mapa com zona destacada."""
    fig, ax = plt.subplots(figsize=figsize, facecolor='none')
    for z in ORDER:
        sub = df[df['zona']==z]
        if z == zona_focus:
            ax.scatter(sub['lon'], sub['lat'], c=ZONE_HEX[z],
                       s=sub['pop']/300, alpha=0.7, linewidths=0, zorder=3, label=z)
        else:
            ax.scatter(sub['lon'], sub['lat'], c='#2a3444',
                       s=4, alpha=0.3, linewidths=0, zorder=1)
    ax.set_facecolor('none')
    ax.tick_params(colors=MIST_HEX, labelsize=7)
    for sp in ax.spines.values(): sp.set_color('#1E283C')
    ax.set_title(f'Localização — {zona_focus}', color='white', fontsize=10,
                 fontweight='bold', pad=5)
    plt.tight_layout(pad=0.3)
    return fig

def make_violin_zona(figsize=(7.5,3.8)):
    fig, ax = plt.subplots(figsize=figsize, facecolor='none')
    data = [df[df['zona']==z]['valor'].values for z in ORDER]
    vp = ax.violinplot(data, positions=range(len(ORDER)), showmedians=True,
                       showextrema=False)
    for i,(body,z) in enumerate(zip(vp['bodies'],ORDER)):
        body.set_facecolor(ZONE_HEX[z]); body.set_alpha(0.65); body.set_edgecolor('none')
    vp['cmedians'].set_color('white'); vp['cmedians'].set_linewidth(1.5)
    ax.set_xticks(range(len(ORDER))); ax.set_xticklabels(ORDER, fontsize=8)
    ax.set_ylabel('Valor Mediano ($)', color=MIST_HEX, fontsize=8)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_: f'${x/1e3:.0f}k'))
    styled_ax(ax, 'Distribuição de Preços por Zona')
    plt.tight_layout(pad=0.3)
    return fig

def make_bar_h(labels, values, colors, title, xlabel, figsize=(5.5,2.8)):
    fig, ax = plt.subplots(figsize=figsize, facecolor='none')
    bars = ax.barh(labels, values, color=colors, height=0.5, edgecolor='none')
    for b, v in zip(bars, values):
        ax.text(v + max(values)*0.01, b.get_y()+b.get_height()/2,
                f'{v:,.0f}' if v>1000 else f'{v:.2f}',
                va='center', color='white', fontsize=8.5, fontweight='bold')
    styled_ax(ax, title)
    ax.set_xlabel(xlabel, color=MIST_HEX, fontsize=8)
    ax.xaxis.grid(True, color='#1E283C', lw=0.5); ax.set_axisbelow(True)
    ax.yaxis.grid(False)
    plt.tight_layout(pad=0.3)
    return fig

def make_scatter_renda_valor(figsize=(6.0,4.2)):
    fig, ax = plt.subplots(figsize=figsize, facecolor='none')
    sample = df.sample(3000, random_state=42)
    for z in ORDER:
        sub = sample[sample['zona']==z]
        ax.scatter(sub['renda']*10000, sub['valor'],
                   c=ZONE_HEX[z], s=8, alpha=0.4, label=z, linewidths=0)
    from scipy import stats as scipy_stats
    m,b,r,_,_ = scipy_stats.linregress(df['renda']*10000, df['valor'])
    xs = np.linspace(df['renda'].min()*10000, df['renda'].max()*10000, 100)
    ax.plot(xs, m*xs+b, color='#F5A623', lw=2, linestyle='--',
            label=f'Tendência (r={r:.2f})')
    ax.legend(fontsize=7, labelcolor='white', facecolor='none',
              edgecolor='none', ncol=2)
    ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_: f'${x/1e3:.0f}k'))
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_: f'${x/1e3:.0f}k'))
    styled_ax(ax, 'Renda Anual × Valor do Imóvel')
    ax.set_xlabel('Renda Anual do Distrito', color=MIST_HEX, fontsize=8)
    ax.set_ylabel('Valor Mediano ($)', color=MIST_HEX, fontsize=8)
    plt.tight_layout(pad=0.3)
    return fig

def make_histo_zona(zona, color, figsize=(4.5,2.8)):
    sub = df[df['zona']==zona]['valor']
    fig, ax = plt.subplots(figsize=figsize, facecolor='none')
    ax.hist(sub, bins=40, color=color, alpha=0.75, edgecolor='none')
    ax.axvline(sub.median(), color='white', lw=1.5, linestyle='--',
               label=f'Mediana ${sub.median()/1e3:.0f}k')
    ax.axvline(sub.mean(),   color='#F5A623', lw=1.5, linestyle=':',
               label=f'Média ${sub.mean()/1e3:.0f}k')
    ax.legend(fontsize=7, labelcolor='white', facecolor='none', edgecolor='none')
    ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_: f'${x/1e3:.0f}k'))
    styled_ax(ax, f'Distribuição de Preços — {zona}')
    ax.set_ylabel('Nº de Distritos', color=MIST_HEX, fontsize=8)
    plt.tight_layout(pad=0.3)
    return fig

def make_score_bar(figsize=(5.5,2.5)):
    scores = {z: ZS.loc[z,'sc'] for z in ORDER}
    fig, ax = plt.subplots(figsize=figsize, facecolor='none')
    zs, vs = zip(*sorted(scores.items(), key=lambda x: -x[1]))
    colors = [ZONE_HEX[z] for z in zs]
    bars = ax.barh(zs, vs, color=colors, height=0.5, edgecolor='none')
    for b,v in zip(bars,vs):
        ax.text(v+0.02, b.get_y()+b.get_height()/2,
                f'{v:.2f}', va='center', color='white', fontsize=9, fontweight='bold')
    styled_ax(ax, 'Score: Renda ÷ (Valor/$100k)')
    ax.xaxis.grid(True, color='#1E283C', lw=0.5); ax.set_axisbelow(True)
    ax.yaxis.grid(False)
    plt.tight_layout(pad=0.3)
    return fig

def make_waterfall_cost(figsize=(6.5,3.5)):
    """Comparação custo mensal estimado por zona."""
    zonas_ord = ['INLAND','<1H OCEAN','NEAR OCEAN','NEAR BAY','ISLAND']
    # Estimativa: prestação 30 anos, 20% entrada, taxa 7% a.a.
    valores = {z: ZS.loc[z,'vm'] for z in zonas_ord}
    rendas_anuais = {z: ZS.loc[z,'rm']*10000 for z in zonas_ord}

    prestacoes = {}
    for z,v in valores.items():
        principal = v * 0.8
        taxa_mensal = 0.07/12
        n = 360
        pmt = principal * (taxa_mensal*(1+taxa_mensal)**n) / ((1+taxa_mensal)**n - 1)
        prestacoes[z] = pmt

    renda_mensal = {z: r/12 for z,r in rendas_anuais.items()}
    comprometimento = {z: prestacoes[z]/renda_mensal[z]*100 for z in zonas_ord}

    fig, axes = plt.subplots(1,2, figsize=figsize, facecolor='none')
    colors = [ZONE_HEX[z] for z in zonas_ord]

    bars1 = axes[0].bar(zonas_ord, [prestacoes[z] for z in zonas_ord],
                        color=colors, edgecolor='none', width=0.55)
    axes[0].axhline(875, color='#F5A623', lw=1.5, linestyle='--',
                    label='Limite saudável\n($875/mês)')
    for b,v in zip(bars1, [prestacoes[z] for z in zonas_ord]):
        axes[0].text(b.get_x()+b.get_width()/2, b.get_height()+10,
                    f'${v:,.0f}', ha='center', color='white', fontsize=7.5, fontweight='bold')
    axes[0].legend(fontsize=7, labelcolor='white', facecolor='none', edgecolor='none')
    styled_ax(axes[0], 'Prestação Mensal Estimada')
    axes[0].set_ylabel('$/mês', color=MIST_HEX, fontsize=8)
    axes[0].tick_params(axis='x', rotation=20)

    bar2_colors = ['#34D399' if v<=30 else '#F5A623' if v<=50 else '#F43F5E'
                   for v in comprometimento.values()]
    bars2 = axes[1].bar(zonas_ord, list(comprometimento.values()),
                        color=bar2_colors, edgecolor='none', width=0.55)
    axes[1].axhline(30, color='#34D399', lw=1.5, linestyle='--', label='Ideal (30%)')
    axes[1].axhline(50, color='#F43F5E', lw=1.5, linestyle='--', label='Crítico (50%)')
    for b,v in zip(bars2, comprometimento.values()):
        axes[1].text(b.get_x()+b.get_width()/2, v+0.5,
                    f'{v:.0f}%', ha='center', color='white', fontsize=7.5, fontweight='bold')
    axes[1].legend(fontsize=7, labelcolor='white', facecolor='none', edgecolor='none')
    styled_ax(axes[1], '% da Renda Comprometida')
    axes[1].set_ylabel('% da Renda Mensal', color=MIST_HEX, fontsize=8)
    axes[1].tick_params(axis='x', rotation=20)

    plt.tight_layout(pad=0.4)
    return fig, prestacoes, comprometimento

# ══════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ══════════════════════════════════════════════════════════════════════
s = slide()
# Gradiente simulado — retângulo dourado lateral grande
rect(s, 0, 0, 5.8, 7.5, BG2)
rect(s, 0, 0, 0.12, 7.5, GOLD)

# Título dramático
txt(s, 'SEU PRÓXIMO', 0.35, 1.1, 5.3, 0.75, size=38, bold=True, color=WHITE)
txt(s, 'LAR NA', 0.35, 1.8, 5.3, 0.75, size=38, bold=True, color=WHITE)
txt(s, 'CALIFÓRNIA', 0.35, 2.5, 5.3, 0.8, size=40, bold=True, color=GOLD)

rect(s, 0.35, 3.45, 4.8, 0.04, GOLD)

txt(s, 'Uma análise orientada por dados para responder\na pergunta que vale milhares de dólares:',
    0.35, 3.6, 5.1, 0.7, size=12, color=MIST)
txt(s, '"Onde comprar ou alugar?"', 0.35, 4.3, 5.1, 0.45,
    size=13, bold=True, italic=True, color=SAND)

txt(s, '20.640 distritos  ·  5 zonas geográficas  ·  Censo 1990',
    0.35, 5.25, 5.3, 0.35, size=9, color=MIST, italic=True)

tag(s, 'DATA STORYTELLING', 0.35, 5.75, bg=GOLD, fg=BG)
tag(s, 'MERCADO IMOBILIÁRIO', 2.55, 5.75, bg=BG3, fg=GOLD)

# Mapa à direita
fig = make_map('valor', 'plasma', 'Mapa de Valor — Califórnia', figsize=(7.0,6.8))
fig2slide(s, fig, 5.9, 0.25, 7.15, 7.0)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 2 — A PERGUNTA
# ══════════════════════════════════════════════════════════════════════
s = slide()
bar_left(s, 0)
eyebrow(s, 'O PONTO DE PARTIDA')
txt(s, 'O mercado imobiliário da', 0.3, 0.55, 9, 0.6, size=28, bold=True, color=WHITE)
txt(s, 'Califórnia é vasto — e confuso.', 0.3, 1.1, 11, 0.6, size=28, bold=True, color=GOLD)
divider(s, 1.8)

txt(s, 'Com 20.640 distritos censitários espalhados por mais de 400 mil km², '
       'a Califórnia tem imóveis que variam de $77.500 a $500.000+ dentro do mesmo estado. '
       'A diferença não é aleatória — ela segue padrões claros que os dados revelam.',
    0.3, 1.95, 12.5, 0.7, size=12, color=MIST)

# 3 perguntas visuais
perguntas = [
    ('🏡', 'Quero comprar', 'Qual zona oferece melhor\ncusto-benefício real?', GOLD),
    ('🔑', 'Quero alugar',  'Onde o aluguel cabe no\norçamento sem apertar?', OCEAN),
    ('📈', 'Quero investir','Qual zona tem maior\npotencial de valorização?', CORAL),
]
for i, (icon, tit, sub, color) in enumerate(perguntas):
    lx = 0.4 + i * 4.3
    rect(s, lx, 2.85, 3.9, 3.6, BG2)
    rect(s, lx, 2.85, 3.9, 0.07, color)
    txt(s, icon, lx+0.15, 2.98, 3.6, 0.65, size=30, align=PP_ALIGN.CENTER)
    txt(s, tit,  lx+0.12, 3.7, 3.6, 0.4, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, sub,  lx+0.12, 4.18, 3.65, 0.7, size=11, color=MIST, align=PP_ALIGN.CENTER)

txt(s, 'Este deck responde às três com dados.', 0.3, 6.6, 12.5, 0.4,
    size=11, italic=True, color=MIST, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 3 — ENTENDENDO AS ZONAS (EDUCATIVO)
# ══════════════════════════════════════════════════════════════════════
s = slide()
bar_left(s, 0, color=OCEAN)
eyebrow(s, 'GUIA DAS ZONAS', color=OCEAN)
txt(s, 'O que significam', 0.3, 0.55, 9, 0.55, size=26, bold=True, color=WHITE)
txt(s, '"INLAND", "NEAR BAY" e as outras zonas?', 0.3, 1.05, 12, 0.55,
    size=26, bold=True, color=OCEAN)
divider(s, 1.7, color=OCEAN)

txt(s, 'O dataset classifica cada distrito conforme sua proximidade geográfica à água. '
       'Entender essa classificação é fundamental para interpretar os preços.',
    0.3, 1.82, 12.5, 0.45, size=11, color=MIST)

zonas_info = [
    ('ISLAND', '🏝️  ISLAND\nIlha',
     'Os 5 distritos das Ilhas Channel, ao largo de Santa Bárbara. '
     'Território isolado, acessado por barco ou avião. Imóveis de luxo extremo, '
     'quase sem mercado de revenda. Uma curiosidade estatística mais do que uma opção real.',
     'Valor mediano: $414.700\n5 distritos apenas'),
    ('NEAR BAY', '🌉  NEAR BAY\nPróximo à Baía',
     'Distritos à beira da Baía de São Francisco — não do Oceano Pacífico. '
     'Inclui São Francisco, Oakland, Berkeley, Palo Alto e cidades do Vale do Silício. '
     'A região tecnológica mais cara dos EUA. Alta densidade, trânsito intenso, salários altíssimos.',
     'Valor mediano: $233.800\n2.290 distritos — SF Bay Area'),
    ('NEAR OCEAN', '🌊  NEAR OCEAN\nPróximo ao Oceano',
     'Distritos com acesso direto ao Oceano Pacífico: Santa Cruz, Monterey, Malibu, '
     'Santa Bárbara, Half Moon Bay. Paisagem deslumbrante, clima ameno o ano todo. '
     'Premium pela natureza, não pelo emprego — ideal para qualidade de vida.',
     'Valor mediano: $229.450\n2.658 distritos — litoral pacífico'),
    ('<1H OCEAN', '🚗  <1H OCEAN\nMenos de 1h do Oceano',
     'Distritos a até 60 minutos de carro da costa, mas sem acesso direto. '
     'A maior zona do dataset. Abrange subúrbios de Los Angeles (Riverside, San Bernardino), '
     'East Bay (Oakland, Fremont), Sacramento e Silicon Valley interior. '
     'O equilíbrio entre acesso e custo — a escolha da maioria das famílias californianas.',
     'Valor mediano: $214.850\n9.136 distritos — a maior zona'),
    ('INLAND', '🏜️  INLAND\nInterior',
     'Distritos a mais de 1 hora do oceano. O coração agrícola da Califórnia: '
     'Vale de San Joaquin, Vale Sacramento, Bakersfield, Fresno, Modesto, Stockton. '
     'Clima extremo (verões de 40°C+), menor IDH, mas o imóvel mais acessível do estado. '
     'O ponto de entrada para quem quer construir patrimônio com renda modesta.',
     'Valor mediano: $108.500\n6.551 distritos — o interior agrícola'),
]

y = 2.4
for zona, header, desc, stats in zonas_info:
    c = rgb(*ZONE_COLOR[zona])
    rect(s, 0.3, y, 12.7, 0.88, BG2)
    rect(s, 0.3, y, 0.07, 0.88, c)
    txt(s, header, 0.5,  y+0.07, 2.2, 0.75, size=9,  bold=True, color=c)
    txt(s, desc,   2.75, y+0.07, 7.3, 0.75, size=8.5, color=MIST)
    txt(s, stats,  10.1, y+0.07, 2.8, 0.75, size=8.5, bold=True, color=WHITE)
    y += 0.97

zone_badge(s, None)
# Override — mostrar todas ativas
rect(s, 0.3, 6.9, 12.7, 0.38, BG3)
zonas_lb = ['Interior','<1h Oceano','Próx. Oceano','Próx. Baía','Ilha']
zonas_k  = ['INLAND','<1H OCEAN','NEAR OCEAN','NEAR BAY','ISLAND']
bw = 12.7/5
for i,(z,lb) in enumerate(zip(zonas_k,zonas_lb)):
    c = rgb(*ZONE_COLOR[z])
    rect(s, 0.3+i*bw, 6.9, bw, 0.38, c)
    txt(s, lb, 0.35+i*bw, 6.96, bw-0.08, 0.28,
        size=8, bold=True, color=BG, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 4 — O PANORAMA GERAL (KPIs + Violin)
# ══════════════════════════════════════════════════════════════════════
s = slide()
bar_left(s, 0, color=CORAL)
eyebrow(s, 'PANORAMA GERAL', color=CORAL)
txt(s, '20.640 distritos.', 0.3, 0.55, 8, 0.55, size=26, bold=True, color=WHITE)
txt(s, 'O que os números revelam?', 0.3, 1.05, 9, 0.5, size=22, bold=True, color=CORAL)
divider(s, 1.68, color=CORAL)

kpi(s, 0.3,  1.85, 2.9, 1.25, '$179.700', 'Valor mediano geral',       GOLD,   '20.640 distritos')
kpi(s, 3.35, 1.85, 2.9, 1.25, '$35.300',  'Renda familiar mediana/ano',SAGE,   '3.53 × $10.000')
kpi(s, 6.4,  1.85, 2.9, 1.25, '$108.500', 'Menor mediana (INLAND)',    CORAL,  'entrada mais acessível')
kpi(s, 9.45, 1.85, 2.9, 1.25, '$414.700', 'Maior mediana (ISLAND)',    RED,    '3.8× o valor do interior')

txt(s, '💡  A amplitude de preços entre a zona mais cara e a mais barata é de 3,8×. '
       'O mesmo estado, com imóveis que variam de $77.500 a $500.000+. '
       'Zona e renda local explicam quase tudo.',
    0.3, 3.25, 12.5, 0.52, size=10.5, italic=True, color=MIST)

fig = make_violin_zona(figsize=(8.5,3.2))
fig2slide(s, fig, 0.3, 3.9, 8.5, 3.2)

# Tabela lateral
rect(s, 9.0, 3.9, 4.05, 3.2, BG2)
txt(s, 'MEDIANA POR ZONA', 9.1, 3.98, 3.8, 0.32, size=9, bold=True, color=GOLD)
headers_t = ['Zona','Mediana','Renda','Score CB']
col_ls = [9.05, 10.25, 11.25, 12.25]
col_ws = [1.15, 0.95, 0.95, 0.75]
for j,(h,l,w) in enumerate(zip(headers_t, col_ls, col_ws)):
    txt(s, h, l, 4.32, w, 0.22, size=7.5, bold=True, color=MIST)
rect(s, 9.05, 4.54, 3.9, 0.025, BG3)
for i,z in enumerate(ORDER):
    yr = 4.6 + i*0.47
    c = rgb(*ZONE_COLOR[z])
    row_bg = BG3 if i%2==0 else BG2
    rect(s, 9.05, yr, 3.9, 0.43, row_bg)
    rect(s, 9.05, yr, 0.05, 0.43, c)
    vals = [z, f"${ZS.loc[z,'vm']/1e3:.0f}k",
            f"{ZS.loc[z,'rm']:.2f}×$10k", f"{ZS.loc[z,'sc']:.2f}"]
    for val,l,w in zip(vals,col_ls,col_ws):
        bold = val==z
        col = c if val==z else WHITE
        txt(s, val, l+0.06, yr+0.07, w, 0.3, size=8, bold=bold, color=col)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 5 — A LEI DA RENDA
# ══════════════════════════════════════════════════════════════════════
s = slide()
bar_left(s, 0, color=GOLD)
eyebrow(s, 'O FATOR INVISÍVEL', color=GOLD)
txt(s, 'A renda do bairro define o preço.', 0.3, 0.55, 10, 0.55, size=24, bold=True, color=WHITE)
txt(s, 'Não o tamanho, não a idade — a renda.', 0.3, 1.05, 10, 0.5, size=20, bold=True, color=GOLD)
divider(s, 1.68)

fig = make_scatter_renda_valor(figsize=(6.5,4.5))
fig2slide(s, fig, 0.3, 1.85, 6.5, 4.5)

# Insight cards
insights = [
    ('Spearman ρ = +0.68', 'Renda × Valor',
     'A correlação mais forte do dataset. IC 95% Bootstrap: [0.67, 0.69]. '
     'Onde a renda é alta, o valor é alto — sempre.', GOLD),
    ('IC 95%: [0.67, 0.69]', 'Intervalo de Confiança',
     'Calculado com 2.000 reamostras Bootstrap. '
     'A correlação é robusta e não é fruto do acaso.', SAGE),
    ('Renda > Localização', 'O que mais importa',
     'Um imóvel no interior com renda alta vale mais '
     'do que um imóvel costeiro em região de baixa renda. '
     'O bairro define o preço.', CORAL),
]
for i,(val,lab,desc,color) in enumerate(insights):
    callout(s, 7.05, 1.85 + i*1.6, 5.9, 1.45, desc, f'{val} — {lab}', color)

txt(s, '📊  O modelo de Gradient Boosting com renda como feature principal atingiu R²=0.83 — '
       'explicando 83% da variação de preços.', 0.3, 6.5, 12.5, 0.45,
    size=10, color=MIST, italic=True)

# ══════════════════════════════════════════════════════════════════════
# HELPER — slide de zona detalhada
# ══════════════════════════════════════════════════════════════════════
def zona_slide(zona, icon, nome_curto, nome_longo, descricao_longa,
               ideal_para, cuidados, cidades, cor):

    s = slide()
    c = rgb(*cor)
    bar_left(s, 0, color=c)
    zone_badge(s, zona)

    tag(s, f'ZONA EM DETALHE', 0.3, 0.22, bg=c, fg=BG)
    txt(s, f'{icon}  {nome_curto}', 0.3, 0.55, 9, 0.6, size=24, bold=True, color=c)
    txt(s, nome_longo, 0.3, 1.12, 9, 0.42, size=14, color=WHITE)
    divider(s, 1.65, color=c)

    # KPIs da zona
    row = ZS.loc[zona]
    kpi(s, 0.3,  1.82, 2.55, 1.1, f"${row['vm']/1e3:.0f}k",  'Valor Mediano', c,     f"P25: ${row['vp25']/1e3:.0f}k / P75: ${row['vp75']/1e3:.0f}k")
    kpi(s, 3.0,  1.82, 2.55, 1.1, f"{row['rm']:.2f}×$10k",   'Renda Mediana', GOLD,  f"${row['rm']*10000:,.0f}/ano")
    kpi(s, 5.7,  1.82, 2.55, 1.1, f"{row['sc']:.2f}",         'Score C/B',     SAGE,  'renda ÷ (valor/$100k)')
    kpi(s, 8.4,  1.82, 2.55, 1.1, f"{int(row['im'])} anos",   'Idade Mediana', MIST,  f"{row['n']:,} distritos")
    kpi(s, 11.1, 1.82, 1.95, 1.1, f"{row['qm']:.1f}",         'Quartos/Dom.',  CORAL, 'média da zona')

    # Descrição
    rect(s, 0.3, 3.1, 6.5, 1.55, BG2)
    rect(s, 0.3, 3.1, 0.07, 1.55, c)
    txt(s, descricao_longa, 0.5, 3.18, 6.1, 1.4, size=9.5, color=MIST)

    # Histograma
    fig = make_histo_zona(zona, '#{:02X}{:02X}{:02X}'.format(*cor), figsize=(4.5,2.8))
    fig2slide(s, fig, 7.0, 3.05, 5.95, 2.6)

    # Mapa da zona
    fig2 = make_zone_scatter(zona, figsize=(4.5,2.5))
    # Não tem mapa — usa cards
    # fig2slide(s, fig2, 7.0, 5.5, 5.95, 1.5)

    # Cards de ideal e cuidados
    y = 4.82
    txt(s, '✅  IDEAL PARA:', 0.35, y, 3.0, 0.3, size=9, bold=True, color=SAGE)
    for item in ideal_para:
        y += 0.32
        txt(s, f'  {item}', 0.35, y, 6.3, 0.28, size=9, color=WHITE)

    y = 4.82
    txt(s, '⚠️  ATENÇÃO:', 6.6, y, 3.0, 0.3, size=9, bold=True, color=CORAL)
    for item in cuidados:
        y += 0.32
        txt(s, f'  {item}', 6.6, y, 6.2, 0.28, size=9, color=MIST)

    # Cidades referência
    rect(s, 0.3, 6.3, 12.7, 0.52, BG3)
    txt(s, '📍  Cidades referência: ', 0.45, 6.38, 2.4, 0.34,
        size=9, bold=True, color=c)
    txt(s, cidades, 2.85, 6.38, 10.0, 0.34, size=9, color=WHITE)

# ══════════════════════════════════════════════════════════════════════
# SLIDES 6–9 — ZONAS DETALHADAS
# ══════════════════════════════════════════════════════════════════════
zona_slide(
    zona='INLAND', icon='🏜️', nome_curto='INLAND — Interior',
    nome_longo='Mais de 1 hora de carro até o oceano. O coração agrícola da Califórnia.',
    descricao_longa=(
        'O INLAND é a Califórnia que os turistas não veem. São os vales agrícolas que '
        'abastecem os EUA — algodão, amêndoas, uvas. As cidades cresceram rápido nos anos '
        '80 e 90 impulsionadas por famílias que não podiam pagar a costa. Hoje, é onde a '
        'classe trabalhadora construiu patrimônio. Verões de 38–42°C são a principal '
        'desvantagem. A infraestrutura melhora a cada década.'
    ),
    ideal_para=[
        '✅  Primeira moradia com renda modesta (< $45k/ano)',
        '✅  Famílias grandes que precisam de espaço por menos',
        '✅  Investidores buscando imóveis abaixo de $150k',
        '✅  Aposentados com renda fixa que querem patrimônio',
    ],
    cuidados=[
        '⚠️  Verões extremamente quentes (38–42°C)',
        '⚠️  Dependência total do carro — sem transporte público',
        '⚠️  Menor valorização histórica vs. zona costeira',
        '⚠️  Mercado de trabalho mais restrito',
    ],
    cidades='Bakersfield · Fresno · Modesto · Stockton · Visalia · Merced · Sacramento (periferia)',
    cor=ZONE_COLOR['INLAND'],
)

zona_slide(
    zona='<1H OCEAN', icon='🚗', nome_curto='<1H OCEAN — Subúrbio Conectado',
    nome_longo='Até 60 minutos do oceano. O endereço da classe média californiana.',
    descricao_longa=(
        'É a zona mais populosa do dataset — 9.136 distritos, 44% do total. '
        'Abrange os subúrbios que cresceram quando SF e LA ficaram inacessíveis. '
        'Riverside e San Bernardino explodiram nos anos 2000. O East Bay virou '
        'alternativa à San Francisco pós-2010. Boa infraestrutura, acesso a empregos '
        'tech por BART ou rodovia, e praias a um fim de semana de distância.'
    ),
    ideal_para=[
        '✅  Famílias que trabalham em SF ou LA mas não podem pagar a cidade',
        '✅  Jovens profissionais de tech em crescimento de carreira',
        '✅  Quem quer escola pública boa + espaço + custo controlado',
        '✅  Primeira compra com renda entre $50–80k/ano',
    ],
    cuidados=[
        '⚠️  Commute longo — 45min a 1h30 por dia',
        '⚠️  Mercado competitivo — demanda alta, estoque baixo',
        '⚠️  Alta variação interna: $100k a $400k+ na mesma zona',
        '⚠️  Risco de congestionamento crônico',
    ],
    cidades='Riverside · San Bernardino · Fremont · Oakland (periferia) · Concord · Stockton · Sacramento',
    cor=ZONE_COLOR['<1H OCEAN'],
)

zona_slide(
    zona='NEAR BAY', icon='🌉', nome_curto='NEAR BAY — Baía de São Francisco',
    nome_longo='Às margens da Baía de SF. O epicentro da tecnologia e do custo de vida americano.',
    descricao_longa=(
        'NEAR BAY não é o Pacífico — é a Baía de São Francisco. Uma distinção importante. '
        'São Francisco, Oakland, Berkeley, Palo Alto, San Jose. O Vale do Silício mora aqui. '
        'A mediana de $233.800 é "barata" para os padrões de 2024, mas em 1990 era premium. '
        'A idade mediana dos imóveis é 39 anos — as casas são antigas, o que encarece manutenção. '
        'O mercado de trabalho é o mais competitivo do planeta.'
    ),
    ideal_para=[
        '✅  Profissionais de tecnologia com salário de $100k+',
        '✅  Investimento de longo prazo (30+ anos de valorização histórica)',
        '✅  Quem prioriza acesso a transporte público (BART)',
        '✅  Perfil urbano — sem dependência de carro',
    ],
    cuidados=[
        '⚠️  Custo de vida altíssimo em tudo (não só o imóvel)',
        '⚠️  Imóveis antigos: manutenção cara e frequente',
        '⚠️  Mercado supercompetitivo — guerras de lances são comuns',
        '⚠️  Risco de bolha — historicamente volátil',
    ],
    cidades='São Francisco · Oakland · Berkeley · Palo Alto · San Jose · Fremont · Richmond',
    cor=ZONE_COLOR['NEAR BAY'],
)

zona_slide(
    zona='NEAR OCEAN', icon='🌊', nome_curto='NEAR OCEAN — Litoral Pacífico',
    nome_longo='Acesso direto ao Oceano Pacífico. Natureza, clima perfeito, ritmo de vida.',
    descricao_longa=(
        'NEAR OCEAN são os distritos que têm o Pacífico como quintal. '
        'Santa Cruz, Monterey, Malibu, Santa Bárbara, Half Moon Bay. '
        'O clima é o melhor da Califórnia: 18–24°C o ano todo, sem extremos. '
        'A valorização é sustentada pela escassez — não há mais litoral para construir. '
        'O perfil de renda é mais baixo que NEAR BAY, mas o preço é quase igual, '
        'porque as pessoas pagam pelo estilo de vida, não pelo mercado de trabalho.'
    ),
    ideal_para=[
        '✅  Quem prioriza qualidade de vida acima de tudo',
        '✅  Trabalho remoto — sem necessidade de ir ao escritório',
        '✅  Aposentados com patrimônio já construído',
        '✅  Investimento em temporada / aluguel por temporada',
    ],
    cuidados=[
        '⚠️  Renda local mais baixa que o preço sugere — crédito exigido',
        '⚠️  Risco de incêndio e erosão costeira em algumas áreas',
        '⚠️  Mercado de trabalho local limitado',
        '⚠️  Custo de seguro imobiliário crescente',
    ],
    cidades='Santa Cruz · Monterey · Carmel · Santa Bárbara · Malibu · Half Moon Bay · Oxnard',
    cor=ZONE_COLOR['NEAR OCEAN'],
)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 10 — ONDE COMPRAR
# ══════════════════════════════════════════════════════════════════════
s = slide()
bar_left(s, 0, color=SAGE)
eyebrow(s, 'DECISÃO DE COMPRA', color=SAGE)
txt(s, 'Onde comprar?', 0.3, 0.55, 8, 0.55, size=26, bold=True, color=WHITE)
txt(s, 'A resposta depende da sua renda e do seu objetivo.', 0.3, 1.08, 10, 0.45,
    size=14, color=SAGE)
divider(s, 1.65, color=SAGE)

fig_score = make_score_bar(figsize=(6.0,3.0))
fig2slide(s, fig_score, 0.3, 1.8, 6.0, 3.2)

fig_w, prestacoes, comprometimento = make_waterfall_cost(figsize=(6.8,3.2))
fig2slide(s, fig_w, 6.5, 1.8, 6.6, 3.2)

# Recomendações
recomendacoes_compra = [
    ('INLAND',    'Renda < $45k/ano',   'Entrada a partir de $22k (20%). Score CB 2.75 — melhor custo-benefício do dataset.',     rgb(*ZONE_COLOR['INLAND'])),
    ('<1H OCEAN', 'Renda $50–80k/ano',  'Equilíbrio real. Prestação ~$1.140/mês. Compromete ~35% da renda — dentro do aceitável.', rgb(*ZONE_COLOR['<1H OCEAN'])),
    ('NEAR BAY',  'Renda $100k+/ano',   'Investimento sólido de longo prazo. Valorização histórica comprovada. Exige fôlego.',      rgb(*ZONE_COLOR['NEAR BAY'])),
]
y = 5.2
for zona, perfil, desc, c in recomendacoes_compra:
    rect(s, 0.3, y, 12.7, 0.62, BG2)
    rect(s, 0.3, y, 0.07, 0.62, c)
    txt(s, zona,   0.5,  y+0.07, 2.0, 0.28, size=10, bold=True, color=c)
    txt(s, perfil, 2.55, y+0.07, 2.4, 0.28, size=9.5, bold=True, color=WHITE)
    txt(s, desc,   5.1,  y+0.07, 7.8, 0.28, size=9, color=MIST)
    y += 0.72

rect(s, 0.3, 7.22, 12.7, 0.22, BG3)
txt(s, '📌  Regra: prestação ≤ 30% da renda bruta mensal. Simular com entrada de 20%, 30 anos, 7% a.a.',
    0.45, 7.25, 12.3, 0.18, size=8.5, color=MIST, italic=True)

zone_badge(s, None)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 11 — ONDE ALUGAR
# ══════════════════════════════════════════════════════════════════════
s = slide()
bar_left(s, 0, color=OCEAN)
eyebrow(s, 'DECISÃO DE ALUGUEL', color=OCEAN)
txt(s, 'Onde alugar?', 0.3, 0.55, 8, 0.55, size=26, bold=True, color=WHITE)
txt(s, 'Custo mensal, liberdade e qualidade de vida como critérios.', 0.3, 1.08, 10, 0.45,
    size=13, color=OCEAN)
divider(s, 1.65, color=OCEAN)

# Estimativa de aluguel: 0.8% do valor/mês (gross rent multiplier típico CA anos 90)
alugueis_est = {z: ZS.loc[z,'vm'] * 0.008 for z in ORDER}
renda_mensal = {z: ZS.loc[z,'rm'] * 10000 / 12 for z in ORDER}
comprometimento_alug = {z: alugueis_est[z]/renda_mensal[z]*100 for z in ORDER}

fig_al, ax = plt.subplots(figsize=(7.0,3.3), facecolor='none')
width = 0.38
xs = np.arange(len(ORDER))
b1 = ax.bar(xs-width/2, [alugueis_est[z] for z in ORDER],
            width=width, color=[ZONE_HEX[z] for z in ORDER], alpha=0.85, label='Aluguel Est. ($/mês)')
b2 = ax.bar(xs+width/2, [renda_mensal[z] for z in ORDER],
            width=width, color=[ZONE_HEX[z] for z in ORDER], alpha=0.4, label='Renda Mensal ($)')
for b,v in list(zip(b1,[alugueis_est[z] for z in ORDER]))+list(zip(b2,[renda_mensal[z] for z in ORDER])):
    ax.text(b.get_x()+b.get_width()/2, b.get_height()+15,
            f'${v:,.0f}', ha='center', color='white', fontsize=7, fontweight='bold')
ax.axhline(875, color='#F5A623', lw=1.5, linestyle='--', label='Limite saudável ($875/mês)')
ax.set_xticks(xs); ax.set_xticklabels(ORDER, fontsize=8)
ax.legend(fontsize=7.5, labelcolor='white', facecolor='none', edgecolor='none')
styled_ax(ax, 'Aluguel Estimado vs. Renda Mensal por Zona')
ax.set_ylabel('$/mês', color=MIST_HEX, fontsize=8)
plt.tight_layout(pad=0.3)
fig2slide(s, fig_al, 0.3, 1.85, 7.0, 3.5)

recomendacoes_alug = [
    ('🏜️  INLAND',    'Melhor custo',    'Aluguel ~$868/mês · 29% da renda', 'Dentro do ideal de 30%. Única zona financeiramente saudável para a renda mediana local.', ZONE_COLOR['INLAND']),
    ('🚗  <1H OCEAN', 'Custo-benefício', 'Aluguel ~$1.719/mês · 44% da renda', 'Fora do ideal, mas aceitável se a renda for > $50k. Acesso real ao litoral.', ZONE_COLOR['<1H OCEAN']),
    ('🌊  NEAR OCEAN', 'Qualidade vida', 'Aluguel ~$1.836/mês · 50% da renda', 'Crítico para a renda local. Considere apenas com renda superior a $55k.', ZONE_COLOR['NEAR OCEAN']),
    ('🌉  NEAR BAY',  'Evitar modesto', 'Aluguel ~$1.870/mês · 49% da renda', 'Quase metade da renda mediana. Sem patrimônio sendo construído. Só com salário tech.', ZONE_COLOR['NEAR BAY']),
]
y = 5.5
for icon_label, subtitulo, dado, desc, cor in recomendacoes_alug:
    c = rgb(*cor)
    rect(s, 7.5, y, 5.65, 0.52, BG2)
    rect(s, 7.5, y, 0.06, 0.52, c)
    txt(s, icon_label, 7.65, y+0.04, 2.1, 0.22, size=8.5, bold=True, color=c)
    txt(s, dado,       9.8,  y+0.04, 3.2, 0.22, size=8,   bold=True, color=WHITE)
    txt(s, desc,       7.65, y+0.26, 5.35, 0.22, size=8,  color=MIST)
    y += 0.62

rect(s, 0.3, 5.5, 7.05, 1.85, BG2)
rect(s, 0.3, 5.5, 0.07, 1.85, GOLD)
txt(s, '💡  A Regra dos 30%', 0.5, 5.58, 6.6, 0.32, size=10, bold=True, color=GOLD)
txt(s,
    'Economistas recomendam gastar no máximo 30% da renda bruta com moradia. '
    'Com a renda mediana do dataset ($35.300/ano = $2.942/mês), o teto saudável '
    'é de $882/mês de aluguel.\n\n'
    'Apenas o INLAND está dentro desse limite. Em todas as outras zonas, '
    'a população local está comprometendo 44–50% da renda — '
    'um sinal de pressão habitacional estrutural.',
    0.5, 5.98, 6.5, 1.2, size=9.5, color=MIST)

zone_badge(s, None)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 12 — PERSONAS
# ══════════════════════════════════════════════════════════════════════
s = slide()
bar_left(s, 0, color=CORAL)
eyebrow(s, 'QUAL É O SEU PERFIL?', color=CORAL)
txt(s, 'Cada perfil tem uma estratégia diferente.', 0.3, 0.55, 10, 0.5,
    size=22, bold=True, color=WHITE)
divider(s, 1.1, color=CORAL)

personas = [
    ('👨‍👩‍👧', 'FAMÍLIA EM CRESCIMENTO',
     '$55–75k/ano · 3+ filhos · escolas',
     [('Zona ideal:', '<1H OCEAN', SAGE),
      ('Regiões:', 'Riverside, Concord, Sacramento', WHITE),
      ('Compra:', '~$215k · prestação ~$1.140/mês', GOLD),
      ('Aluguel:', '~$1.720/mês · 37% da renda', OCEAN),
      ('Fator chave:', 'Qualidade de escola pública', MIST)],
     rgb(*ZONE_COLOR['<1H OCEAN'])),

    ('💻', 'JOVEM PROFISSIONAL TECH',
     '$90–130k/ano · solteiro/casal · mobilidade',
     [('Zona ideal:', 'NEAR BAY ou <1H OCEAN', SAGE),
      ('Regiões:', 'East Bay, San Jose, Fremont', WHITE),
      ('Compra:', '~$235k · viável com essa renda', GOLD),
      ('Aluguel:', '~$1.870/mês · 26% da renda', OCEAN),
      ('Fator chave:', 'Acesso a BART/metrô', MIST)],
     rgb(*ZONE_COLOR['NEAR BAY'])),

    ('🌿', 'APOSENTADO / RENDA FIXA',
     '$28–40k/ano · conforto · custo baixo',
     [('Zona ideal:', 'INLAND', SAGE),
      ('Regiões:', 'Bakersfield, Fresno, Merced', WHITE),
      ('Compra:', '~$108k · prestação ~$575/mês', GOLD),
      ('Aluguel:', '~$868/mês · 30% da renda', OCEAN),
      ('Fator chave:', 'Baixo custo de vida geral', MIST)],
     rgb(*ZONE_COLOR['INLAND'])),

    ('📊', 'INVESTIDOR IMOBILIÁRIO',
     '$80k+/ano · retorno e valorização',
     [('Zona ideal:', 'NEAR BAY (longo prazo)', SAGE),
      ('Alternativa:', '<1H OCEAN (custo de entrada)', WHITE),
      ('ROI estimado:', '5–7% ao ano histórico em SF', GOLD),
      ('Estratégia:', 'Compra → aluguel → valorização', OCEAN),
      ('Fator chave:', 'Demanda de tech sustentada', MIST)],
     rgb(*ZONE_COLOR['NEAR BAY'])),
]

positions_p = [(0.3,1.3),(3.55,1.3),(6.8,1.3),(10.05,1.3)]
for (icon,nome,renda,items,c), (lx,ty) in zip(personas, positions_p):
    rect(s, lx, ty, 3.1, 5.85, BG2)
    rect(s, lx, ty, 3.1, 0.08, c)
    txt(s, icon, lx+0.12, ty+0.12, 2.85, 0.55, size=28, align=PP_ALIGN.CENTER)
    txt(s, nome, lx+0.1, ty+0.7, 2.9, 0.38, size=10, bold=True, color=c, align=PP_ALIGN.CENTER)
    txt(s, renda, lx+0.1, ty+1.05, 2.9, 0.28, size=7.5, color=MIST, italic=True, align=PP_ALIGN.CENTER)
    rect(s, lx+0.15, ty+1.38, 2.8, 0.025, BG3)
    yi = ty + 1.52
    for label, value, vc in items:
        txt(s, label, lx+0.15, yi, 1.15, 0.3, size=8, bold=True, color=MIST)
        txt(s, value, lx+1.32, yi, 1.65, 0.3, size=8, color=vc)
        yi += 0.38

# ══════════════════════════════════════════════════════════════════════
# SLIDE 13 — RANKING FINAL
# ══════════════════════════════════════════════════════════════════════
s = slide()
bar_left(s, 0, color=GOLD)
eyebrow(s, 'RANKING FINAL', color=GOLD)
txt(s, 'As 5 zonas ranqueadas.', 0.3, 0.55, 8, 0.55, size=24, bold=True, color=WHITE)
txt(s, 'Critérios: valor, custo-benefício, comprometimento de renda, qualidade de vida, valorização.',
    0.3, 1.08, 12.5, 0.42, size=11, color=MIST)
divider(s, 1.62)

ranking = [
    ('#1', '<1H OCEAN',  'COMPRAR & ALUGAR',
     'O equilíbrio perfeito para a maioria das famílias californianas. Acesso ao litoral, '
     'infraestrutura de qualidade, mercado de trabalho diverso. Maior zona do dataset.',
     '★★★★★', '★★★★', rgb(*ZONE_COLOR['<1H OCEAN'])),
    ('#2', 'INLAND',     'COMPRAR (1ª moradia)',
     'Melhor score custo-benefício do dataset. Único onde a renda mediana local sustenta '
     'confortavelmente o custo de moradia. A porta de entrada para o patrimônio.',
     '★★★★★', '★★★', rgb(*ZONE_COLOR['INLAND'])),
    ('#3', 'NEAR BAY',   'INVESTIR (longo prazo)',
     'Valorização histórica mais consistente. Demanda estrutural de tech. '
     'Exige renda alta, mas o retorno justifica para investidores pacientes.',
     '★★★', '★★★★★', rgb(*ZONE_COLOR['NEAR BAY'])),
    ('#4', 'NEAR OCEAN', 'ALUGAR (qualidade de vida)',
     'O litoral acessível. Melhor clima, menor densidade que Bay Area. '
     'Custo acima do ideal, mas justificado pelo estilo de vida — especialmente para remotos.',
     '★★★', '★★★★', rgb(*ZONE_COLOR['NEAR OCEAN'])),
    ('#5', 'ISLAND',     'EVITAR (nicho extremo)',
     'Luxo absoluto sem liquidez. Cinco distritos apenas. Sem mercado de trabalho, '
     'sem infraestrutura. Exclusivamente para patrimônio muito elevado.',
     '★', '★★', rgb(*ZONE_COLOR['ISLAND'])),
]

headers_r = ['Rank','Zona','Objetivo','Por quê?','Acessib.','Valoriz.']
col_l_r = [0.3, 0.95, 2.5,  4.2,  11.2, 12.3]
col_w_r = [0.6, 1.5,  1.65, 7.0,  1.05, 0.95]

rect(s, 0.3, 1.78, 12.7, 0.38, BG3)
for h,l,w in zip(headers_r, col_l_r, col_w_r):
    txt(s, h, l+0.05, 1.83, w, 0.28, size=8.5, bold=True, color=GOLD)

for i,(rank, zona, obj, desc, acesso, valor_v, c) in enumerate(ranking):
    y = 2.24 + i*1.02
    bg = BG2 if i%2==0 else rgb(0x16,0x22,0x34)
    rect(s, 0.3, y, 12.7, 0.96, bg)
    rect(s, 0.3, y, 0.07, 0.96, c)
    data_r = [rank, zona, obj, desc, acesso, valor_v]
    for val,l,w in zip(data_r, col_l_r, col_w_r):
        bold = val in [rank, zona]
        col = c if val in [rank, zona] else (WHITE if val==obj else MIST)
        sz = 15 if val==rank else (10 if val==zona else 8.5)
        txt(s, val, l+0.06, y+0.06, w-0.08, 0.84, size=sz, bold=bold, color=col)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 14 — CONCLUSÃO
# ══════════════════════════════════════════════════════════════════════
s = slide()
rect(s, 0, 0, 13.33, 7.5, BG2)
# Barra dourada topo e base
rect(s, 0, 0,    13.33, 0.12, GOLD)
rect(s, 0, 7.38, 13.33, 0.12, GOLD)

txt(s, 'TRÊS VERDADES QUE OS DADOS PROVAM', 0.4, 0.3, 12.5, 0.6,
    size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txt(s, 'sobre o mercado imobiliário da Califórnia', 0.4, 0.85, 12.5, 0.4,
    size=13, color=MIST, italic=True, align=PP_ALIGN.CENTER)

verdades = [
    ('01', 'A RENDA DO BAIRRO\nDEFINE O PREÇO',
     'Não o tamanho, não a vista, não a proximidade ao oceano de forma isolada. '
     'O Spearman ρ = 0.68 com IC [0.67–0.69] é o número mais importante deste deck. '
     'Antes de olhar o imóvel, pesquise a renda mediana do distrito.',
     GOLD),
    ('02', 'O INLAND É O\nMELHOR NEGÓCIO',
     'Score custo-benefício de 2.75 — o mais alto do dataset. '
     'Único onde a renda local sustenta o custo de moradia dentro dos 30% recomendados. '
     'Ignorado por preconceito estético, não por dados.',
     rgb(*ZONE_COLOR['INLAND'])),
    ('03', 'O <1H OCEAN É\nO EQUILÍBRIO REAL',
     'A maior zona, a mais diversa, a mais conectada. Para 44% dos distritos da Califórnia, '
     'esse é o meio-termo entre sonho e viabilidade. O endereço da classe média que '
     'conseguiu encontrar o ponto de equilíbrio.',
     rgb(*ZONE_COLOR['<1H OCEAN'])),
]

for i, (num, titulo, corpo, c) in enumerate(verdades):
    lx = 0.4 + i * 4.3
    rect(s, lx, 1.45, 3.95, 5.2, BG)
    rect(s, lx, 1.45, 3.95, 0.08, c)
    txt(s, num,    lx+0.12, 1.58, 3.7, 0.6, size=32, bold=True, color=c)
    txt(s, titulo, lx+0.12, 2.2,  3.7, 0.75, size=13, bold=True, color=WHITE)
    rect(s, lx+0.12, 3.0, 3.6, 0.03, BG3)
    txt(s, corpo,  lx+0.12, 3.15, 3.7, 2.9, size=9.5, color=MIST)

rect(s, 0.4, 6.8, 12.5, 0.45, BG3)
txt(s, '📂  Análise completa + código:  github.com/LuissFellipe/houses-prices-analysis',
    0.55, 6.87, 12.2, 0.32, size=9.5, color=MIST, italic=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════
# SALVAR
# ══════════════════════════════════════════════════════════════════════
out = Path('presentation') / 'california_housing_v2.pptx'
out.parent.mkdir(exist_ok=True)
prs.save(out)
print(f"✅  Salvo: {out}")
print(f"   Slides: {len(prs.slides)}")
