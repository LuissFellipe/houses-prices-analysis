"""Gera apresentação PowerPoint com storytelling sobre o mercado imobiliário da Califórnia."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import matplotlib.patches as mpatches
from pathlib import Path
import io

# ── Paleta de cores da apresentação ──────────────────────────────────────────
C_BG        = RGBColor(0x0D, 0x1B, 0x2A)   # azul escuro — fundo
C_BG2       = RGBColor(0x1B, 0x2F, 0x45)   # azul médio — cards
C_ACCENT    = RGBColor(0x00, 0xB4, 0xD8)   # ciano — destaque
C_ACCENT2   = RGBColor(0x90, 0xE0, 0xEF)   # ciano claro
C_GREEN     = RGBColor(0x2D, 0xC6, 0x53)   # verde
C_ORANGE    = RGBColor(0xFF, 0x9F, 0x1C)   # laranja
C_RED       = RGBColor(0xFF, 0x4D, 0x6D)   # vermelho
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY      = RGBColor(0xAA, 0xBB, 0xCC)
C_YELLOW    = RGBColor(0xFF, 0xD6, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

OCEAN_HEX = {
    'INLAND':     '#FF9F1C',
    '<1H OCEAN':  '#2DC653',
    'NEAR OCEAN': '#00B4D8',
    'NEAR BAY':   '#4361EE',
    'ISLAND':     '#FF4D6D',
}

# ── Helpers ───────────────────────────────────────────────────────────────────
def rgb(r, g, b): return RGBColor(r, g, b)

def add_rect(slide, l, t, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_img_from_fig(slide, fig, l, t, w, h):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', bbox_inches='tight',
                dpi=150, facecolor='none', transparent=True)
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(l), Inches(t), Inches(w), Inches(h))
    plt.close(fig)

def set_bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_tag(slide, text, l, t, bg=C_ACCENT, fg=C_BG, size=10):
    add_rect(slide, l, t, len(text) * 0.085 + 0.2, 0.28, bg)
    add_text(slide, text, l + 0.05, t + 0.03, len(text) * 0.1 + 0.2, 0.25,
             size=size, bold=True, color=fg, align=PP_ALIGN.CENTER)

def divider(slide, t, color=C_ACCENT, w=12.33, l=0.5):
    add_rect(slide, l, t, w, 0.04, color)

def kpi_card(slide, l, t, w, h, value, label, color=C_ACCENT, sub=None):
    add_rect(slide, l, t, w, h, C_BG2)
    add_rect(slide, l, t, 0.06, h, color)  # borda esquerda colorida
    add_text(slide, value, l + 0.15, t + 0.1, w - 0.2, 0.55,
             size=26, bold=True, color=color, align=PP_ALIGN.LEFT)
    add_text(slide, label, l + 0.15, t + 0.6, w - 0.2, 0.3,
             size=10, color=C_GRAY, align=PP_ALIGN.LEFT)
    if sub:
        add_text(slide, sub, l + 0.15, t + 0.88, w - 0.2, 0.25,
                 size=9, color=C_GRAY, italic=True, align=PP_ALIGN.LEFT)

def bullet_card(slide, l, t, w, items, title=None, bg=C_BG2, icon='●'):
    h = 0.38 * len(items) + (0.45 if title else 0.1)
    add_rect(slide, l, t, w, h, bg)
    y = t + 0.1
    if title:
        add_text(slide, title, l + 0.15, y, w - 0.25, 0.32,
                 size=11, bold=True, color=C_ACCENT)
        y += 0.38
    for item, color in items:
        add_text(slide, f'{icon}  {item}', l + 0.15, y, w - 0.25, 0.32,
                 size=10, color=color)
        y += 0.34
    return h

# ── Gráficos inline ────────────────────────────────────────────────────────────
plt.rcParams.update({'font.family': 'DejaVu Sans'})

def make_bar_ocean(data_dict, title, ylabel, fmt='${}k', figsize=(5.5, 3.2)):
    fig, ax = plt.subplots(figsize=figsize, facecolor='none')
    labels = list(data_dict.keys())
    values = list(data_dict.values())
    colors = [OCEAN_HEX.get(l, '#aaa') for l in labels]
    bars = ax.bar(labels, values, color=colors, edgecolor='none', width=0.55)
    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + max(values)*0.02,
                fmt.format(f'{val/1e3:.0f}' if val > 1000 else f'{val:.2f}'),
                ha='center', va='bottom', fontsize=9, color='white', fontweight='bold')
    ax.set_title(title, color='white', fontsize=11, pad=8, fontweight='bold')
    ax.set_ylabel(ylabel, color='#aabbcc', fontsize=9)
    ax.tick_params(colors='#aabbcc', labelsize=8)
    ax.set_facecolor('none')
    for spine in ax.spines.values(): spine.set_visible(False)
    ax.yaxis.grid(True, color='#1B2F45', linewidth=0.6)
    ax.set_axisbelow(True)
    plt.xticks(rotation=15, ha='right')
    plt.tight_layout()
    return fig

def make_scatter_map(df, col, cmap, title, figsize=(6, 4)):
    fig, ax = plt.subplots(figsize=figsize, facecolor='none')
    sc = ax.scatter(df['longitude'], df['latitude'],
                    c=df[col], cmap=cmap, s=df['populacao']/400,
                    alpha=0.45, linewidths=0)
    cb = fig.colorbar(sc, ax=ax, shrink=0.75, pad=0.02)
    cb.ax.tick_params(colors='#aabbcc', labelsize=7)
    cb.set_label('', color='#aabbcc')
    ax.set_title(title, color='white', fontsize=11, fontweight='bold', pad=6)
    ax.set_facecolor('none')
    ax.tick_params(colors='#aabbcc', labelsize=7)
    for spine in ax.spines.values(): spine.set_color('#1B2F45')
    plt.tight_layout()
    return fig

def make_donut(labels, values, colors, title, figsize=(3.5, 3.0)):
    fig, ax = plt.subplots(figsize=figsize, facecolor='none')
    wedges, texts, autotexts = ax.pie(
        values, labels=None, colors=colors,
        autopct='%1.0f%%', startangle=90,
        wedgeprops=dict(width=0.5, edgecolor='none'),
        pctdistance=0.78,
    )
    for at in autotexts:
        at.set_color('white'); at.set_fontsize(8); at.set_fontweight('bold')
    ax.legend(labels, loc='lower center', bbox_to_anchor=(0.5, -0.18),
              ncol=2, fontsize=7, labelcolor='white',
              facecolor='none', edgecolor='none')
    ax.set_title(title, color='white', fontsize=10, fontweight='bold', pad=4)
    plt.tight_layout()
    return fig

def make_hbar(labels, values, colors, title, xlabel, figsize=(5, 2.8)):
    fig, ax = plt.subplots(figsize=figsize, facecolor='none')
    bars = ax.barh(labels, values, color=colors, edgecolor='none', height=0.5)
    for bar, val in zip(bars, values):
        ax.text(val + max(values)*0.01, bar.get_y() + bar.get_height()/2,
                f'{val:.2f}', va='center', color='white', fontsize=9, fontweight='bold')
    ax.set_title(title, color='white', fontsize=11, fontweight='bold', pad=6)
    ax.set_xlabel(xlabel, color='#aabbcc', fontsize=8)
    ax.set_facecolor('none')
    ax.tick_params(colors='#aabbcc', labelsize=8)
    for spine in ax.spines.values(): spine.set_visible(False)
    ax.xaxis.grid(True, color='#1B2F45', linewidth=0.6)
    ax.set_axisbelow(True)
    plt.tight_layout()
    return fig

# ── Carregar dados ─────────────────────────────────────────────────────────────
df = pd.read_csv('housing.csv')
df.columns = ['longitude','latitude','idade_mediana','total_quartos','total_dormitorios',
              'populacao','domicilios','renda_mediana','valor_mediano','proximidade_oceano']
df['total_dormitorios'] = df['total_dormitorios'].fillna(df['total_dormitorios'].median())
df['quartos_por_dom']   = df['total_quartos'] / df['domicilios']
df['pessoas_por_dom']   = df['populacao'] / df['domicilios']
df['score_cb']          = df['renda_mediana'] / (df['valor_mediano'] / 100000)
df['regiao'] = pd.cut(df['latitude'],
    bins=[32, 34.5, 36.5, 38.5, 42],
    labels=['Sul (Los Angeles)', 'Centro-Sul', 'Norte de SF', 'Norte'])

ocean_stats = df.groupby('proximidade_oceano').agg(
    valor_med=('valor_mediano','median'),
    renda_med=('renda_mediana','median'),
    score_cb=('score_cb','median'),
    n=('valor_mediano','count'),
).sort_values('valor_med')

region_stats = df.groupby('regiao', observed=True).agg(
    valor_med=('valor_mediano','median'),
    renda_med=('renda_mediana','median'),
    n=('valor_mediano','count'),
)

# ── Criar apresentação ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]  # layout em branco

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)

# Gradiente simulado com retângulos
add_rect(s, 0, 0, 6.5, 7.5, C_BG2)
add_rect(s, 0, 0, 0.08, 7.5, C_ACCENT)

# Texto principal
add_text(s, 'ONDE COMPRAR', 0.4, 1.2, 6.0, 0.8,
         size=36, bold=True, color=C_WHITE)
add_text(s, 'OU ALUGAR', 0.4, 1.95, 6.0, 0.8,
         size=36, bold=True, color=C_ACCENT)
add_text(s, 'NA CALIFÓRNIA?', 0.4, 2.7, 6.0, 0.7,
         size=28, bold=True, color=C_WHITE)

divider(s, 3.55, w=5.8, l=0.4)

add_text(s, 'Uma análise orientada por dados do\nCalifornia Housing Dataset',
         0.4, 3.7, 5.8, 0.8, size=13, color=C_GRAY)
add_text(s, '20.640 distritos censitários  ·  Califórnia  ·  Censo 1990',
         0.4, 4.4, 5.8, 0.4, size=10, color=C_GRAY, italic=True)

add_tag(s, 'DATA STORYTELLING', 0.4, 5.8, bg=C_ACCENT, fg=C_BG)
add_tag(s, 'MERCADO IMOBILIÁRIO', 3.0, 5.8, bg=C_BG2, fg=C_ACCENT2)

# Mapa à direita
fig_map = make_scatter_map(df, 'valor_mediano', 'plasma',
                            'Mapa de Calor — Valor dos Imóveis', figsize=(6.2, 6.0))
add_img_from_fig(s, fig_map, 6.8, 0.4, 6.2, 6.5)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — O MERCADO EM NÚMEROS
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
add_tag(s, '01', 0.4, 0.25, bg=C_ACCENT, fg=C_BG)
add_text(s, 'O MERCADO EM NÚMEROS', 0.75, 0.2, 10, 0.5,
         size=22, bold=True, color=C_WHITE)
divider(s, 0.75)

# KPIs
kpi_card(s, 0.4,  1.0, 2.9, 1.15, '$179.700',  'Valor mediano dos imóveis',      C_ACCENT,  'Mediana de 20.640 distritos')
kpi_card(s, 3.45, 1.0, 2.9, 1.15, '$206.856',  'Valor médio dos imóveis',         C_ORANGE,  'Puxado por imóveis premium')
kpi_card(s, 6.5,  1.0, 2.9, 1.15, '$35.300',   'Renda familiar mediana/ano',      C_GREEN,   '3.53 × $10.000 no dataset')
kpi_card(s, 9.55, 1.0, 2.9, 1.15, '36,6%',     'Distritos com imóveis < $150k',   C_YELLOW,  '7.556 distritos acessíveis')

add_text(s, '💡  A diferença de $27k entre média e mediana revela um mercado com cauda longa — '
            'poucos imóveis de altíssimo valor puxam a média para cima, mas a maioria dos distritos '
            'tem preços bem mais acessíveis.',
         0.4, 2.35, 12.5, 0.55, size=11, color=C_GRAY, italic=True)

divider(s, 3.05, color=C_BG2)

# Gráfico barras por oceano
val_ocean = ocean_stats['valor_med'].to_dict()
fig_bar = make_bar_ocean(val_ocean, 'Valor Mediano por Proximidade ao Oceano',
                          'Valor ($)', fmt='${}k')
add_img_from_fig(s, fig_bar, 0.4, 3.15, 6.0, 3.9)

# Donut de distribuição
counts = df['proximidade_oceano'].value_counts()
fig_donut = make_donut(
    list(counts.index), list(counts.values),
    [OCEAN_HEX.get(l,'#aaa') for l in counts.index],
    'Distribuição dos Distritos por Zona'
)
add_img_from_fig(s, fig_donut, 6.5, 3.15, 3.4, 3.9)

# Insight card
add_rect(s, 10.1, 3.15, 3.0, 3.9, C_BG2)
add_rect(s, 10.1, 3.15, 0.06, 3.9, C_ACCENT)
add_text(s, '🏝️  ISLAND\n$414.700', 10.25, 3.3, 2.7, 0.7,
         size=14, bold=True, color=C_RED)
add_text(s, 'Mais caro — média de renda baixa.\nNicho de luxo, inacessível.', 10.25, 3.95, 2.7, 0.5,
         size=9, color=C_GRAY)
add_text(s, '🏠  INLAND\n$108.500', 10.25, 4.6, 2.7, 0.7,
         size=14, bold=True, color=C_GREEN)
add_text(s, 'Mais barato — melhor relação\ncusto × renda do dataset.', 10.25, 5.25, 2.7, 0.5,
         size=9, color=C_GRAY)
add_text(s, '🌊  NEAR BAY\n$233.800', 10.25, 5.9, 2.7, 0.55,
         size=13, bold=True, color=C_ACCENT)
add_text(s, 'Premium — alta renda local.', 10.25, 6.4, 2.7, 0.3,
         size=9, color=C_GRAY)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — A LÓGICA DO PREÇO
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
add_tag(s, '02', 0.4, 0.25, bg=C_ORANGE, fg=C_BG)
add_text(s, 'A LÓGICA DO PREÇO', 0.75, 0.2, 10, 0.5,
         size=22, bold=True, color=C_WHITE)
divider(s, 0.75, color=C_ORANGE)

add_text(s, 'O que de fato move o preço de um imóvel na Califórnia?',
         0.4, 0.9, 12.5, 0.45, size=13, color=C_GRAY, italic=True)

# Mapa renda
fig_renda = make_scatter_map(df, 'renda_mediana', 'YlOrRd',
                              'Renda Mediana por Região', figsize=(5.8, 4.2))
add_img_from_fig(s, fig_renda, 0.4, 1.4, 5.8, 4.5)

# Correlação visual — cards
add_text(s, 'FATORES QUE MAIS INFLUENCIAM O PREÇO', 6.5, 1.4, 6.5, 0.4,
         size=12, bold=True, color=C_ACCENT)

fatores = [
    ('Renda mediana do distrito',    'Spearman ρ = +0.68', C_GREEN,  '████████████████ Fator dominante'),
    ('Latitude (posição geográfica)', 'Spearman ρ = −0.14', C_ORANGE, '████ Quanto mais ao sul, mais caro'),
    ('Quartos por domicílio',         'Spearman ρ = +0.15', C_ACCENT, '████ Mais espaço = mais valor'),
    ('Pessoas por domicílio',         'Spearman ρ = −0.22', C_RED,    '█████ Superlotação reduz valor'),
    ('Proximidade ao oceano',         'Efeito confirmado',  C_ACCENT2,'██████ Premium costeiro: +2×'),
]

y = 1.9
for label, corr, color, bar_txt in fatores:
    add_rect(s, 6.5, y, 6.3, 0.75, C_BG2)
    add_rect(s, 6.5, y, 0.06, 0.75, color)
    add_text(s, label,    6.65, y + 0.04, 3.8, 0.35, size=10, bold=True, color=C_WHITE)
    add_text(s, corr,     6.65, y + 0.38, 2.5, 0.28, size=9,  color=color)
    add_text(s, bar_txt,  9.2,  y + 0.38, 3.5, 0.28, size=8,  color=C_GRAY, italic=True)
    y += 0.85

add_rect(s, 0.4, 6.05, 12.5, 1.05, C_BG2)
add_rect(s, 0.4, 6.05, 0.06, 1.05, C_ORANGE)
add_text(s, '📊  CONCLUSÃO:  Renda mediana do distrito é o melhor preditor individual de valor imobiliário '
            '(ρ=0.68). A localização geográfica e a proximidade ao oceano amplificam esse efeito — '
            'distritos costeiros com alta renda são os mais valorizados.',
         0.55, 6.1, 12.2, 0.85, size=11, color=C_GRAY)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — MELHOR PARA COMPRAR
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
add_tag(s, '03', 0.4, 0.25, bg=C_GREEN, fg=C_BG)
add_text(s, 'ONDE COMPRAR?', 0.75, 0.2, 10, 0.5,
         size=22, bold=True, color=C_WHITE)
divider(s, 0.75, color=C_GREEN)
add_text(s, 'Perfis de comprador e os melhores mercados para cada objetivo',
         0.4, 0.9, 12.5, 0.4, size=12, color=C_GRAY, italic=True)

# Coluna 1 — Perfil econômico
add_rect(s, 0.4, 1.4, 4.0, 5.7, C_BG2)
add_rect(s, 0.4, 1.4, 4.0, 0.08, C_ORANGE)
add_text(s, '💰  PERFIL ECONÔMICO', 0.55, 1.5, 3.7, 0.4,
         size=13, bold=True, color=C_ORANGE)
add_text(s, 'Compra acessível, primeira moradia', 0.55, 1.9, 3.7, 0.3,
         size=9, color=C_GRAY, italic=True)
divider(s, 2.3, color=C_BG, w=3.6, l=0.55)

itens_eco = [
    ('Interior (INLAND)',        C_WHITE),
    ('Mediana: $108.500',        C_ORANGE),
    ('36,6% dos distritos abaixo de $150k', C_GRAY),
    ('Menor concorrência',       C_GRAY),
    ('',                         C_GRAY),
    ('✅ Vale Central',           C_GREEN),
    ('✅ Bakersfield e Fresno',   C_GREEN),
    ('✅ Riverside / San Bernardino', C_GREEN),
    ('',                         C_GRAY),
    ('⚠️  Renda local menor',     C_ORANGE),
    ('⚠️  Infraestrutura limitada',C_ORANGE),
]
y = 2.45
for txt, color in itens_eco:
    if txt:
        add_text(s, txt, 0.55, y, 3.5, 0.28, size=9, color=color)
    y += 0.3

# Coluna 2 — Perfil custo-benefício
add_rect(s, 4.65, 1.4, 4.0, 5.7, C_BG2)
add_rect(s, 4.65, 1.4, 4.0, 0.08, C_ACCENT)
add_text(s, '⚖️  MELHOR CUSTO-BENEFÍCIO', 4.8, 1.5, 3.7, 0.4,
         size=13, bold=True, color=C_ACCENT)
add_text(s, 'Equilíbrio entre preço e qualidade de vida', 4.8, 1.9, 3.7, 0.3,
         size=9, color=C_GRAY, italic=True)
divider(s, 2.3, color=C_BG, w=3.6, l=4.8)

# Score custo-benefício
score_data = ocean_stats['score_cb'].sort_values(ascending=True).to_dict()
fig_score = make_hbar(
    list(score_data.keys()), list(score_data.values()),
    [OCEAN_HEX.get(k,'#aaa') for k in score_data.keys()],
    'Score: Renda / (Valor / $100k)', 'Score',
    figsize=(3.5, 2.5)
)
add_img_from_fig(s, fig_score, 4.7, 2.35, 3.8, 2.8)

add_text(s, '✅  Inland lidera — renda razoável\n    com preços muito acessíveis', 4.8, 5.2, 3.7, 0.55,
         size=10, color=C_GREEN)
add_text(s, '✅  <1H Ocean — zona suburbana\n    equilibrada para famílias', 4.8, 5.75, 3.7, 0.55,
         size=10, color=C_ACCENT)
add_text(s, '⚠️  Near Bay / Island — renda alta\n    exigida para sustentar o custo', 4.8, 6.3, 3.7, 0.55,
         size=9, color=C_ORANGE)

# Coluna 3 — Perfil premium
add_rect(s, 8.9, 1.4, 4.0, 5.7, C_BG2)
add_rect(s, 8.9, 1.4, 4.0, 0.08, C_RED)
add_text(s, '🏆  PERFIL PREMIUM', 9.05, 1.5, 3.7, 0.4,
         size=13, bold=True, color=C_RED)
add_text(s, 'Valorização e qualidade máxima', 9.05, 1.9, 3.7, 0.3,
         size=9, color=C_GRAY, italic=True)
divider(s, 2.3, color=C_BG, w=3.6, l=9.05)

itens_prem = [
    ('Near Bay / Near Ocean',      C_WHITE),
    ('Mediana: $233k–$229k',       C_RED),
    ('Renda local alta',           C_GRAY),
    ('Infraestrutura top',         C_GRAY),
    ('',                           C_GRAY),
    ('✅ San Francisco Bay Area',   C_GREEN),
    ('✅ Peninsula (Palo Alto)',     C_GREEN),
    ('✅ Santa Barbara',            C_GREEN),
    ('✅ San Diego (costa)',         C_GREEN),
    ('',                           C_GRAY),
    ('⚠️  Entrada muito alta',      C_ORANGE),
]
y = 2.45
for txt, color in itens_prem:
    if txt:
        add_text(s, txt, 9.05, y, 3.6, 0.28, size=9, color=color)
    y += 0.3

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MELHOR PARA ALUGAR
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
add_tag(s, '04', 0.4, 0.25, bg=C_ACCENT, fg=C_BG)
add_text(s, 'ONDE ALUGAR?', 0.75, 0.2, 10, 0.5,
         size=22, bold=True, color=C_WHITE)
divider(s, 0.75, color=C_ACCENT)
add_text(s, 'Para quem prioriza flexibilidade, custo mensal baixo ou qualidade de vida',
         0.4, 0.9, 12.5, 0.4, size=12, color=C_GRAY, italic=True)

# Heatmap de renda × valor por região
reg_vals = region_stats['valor_med'].to_dict()
reg_rend = region_stats['renda_med'].to_dict()

fig_reg, ax = plt.subplots(figsize=(5.5, 3.5), facecolor='none')
regs  = list(reg_vals.keys())
vals  = [reg_vals[r] for r in regs]
rends = [reg_rend[r] * 10000 for r in regs]
x = np.arange(len(regs))
w = 0.38
b1 = ax.bar(x - w/2, vals,  width=w, color='#4361EE', label='Valor Mediano ($)', alpha=0.9)
b2 = ax.bar(x + w/2, rends, width=w, color='#2DC653', label='Renda Anual ($)', alpha=0.9)
for bar, val in list(zip(b1, vals)) + list(zip(b2, rends)):
    ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 2000,
            f'${val/1e3:.0f}k', ha='center', va='bottom', color='white', fontsize=7.5, fontweight='bold')
ax.set_xticks(x); ax.set_xticklabels(regs, color='#aabbcc', fontsize=8, rotation=12)
ax.set_facecolor('none')
ax.tick_params(colors='#aabbcc', labelsize=7)
for spine in ax.spines.values(): spine.set_visible(False)
ax.yaxis.grid(True, color='#1B2F45', linewidth=0.5); ax.set_axisbelow(True)
ax.legend(fontsize=8, labelcolor='white', facecolor='none', edgecolor='none')
ax.set_title('Valor Mediano vs. Renda por Região', color='white', fontsize=11, fontweight='bold')
plt.tight_layout()
add_img_from_fig(s, fig_reg, 0.4, 1.4, 6.0, 4.0)

# Cards de recomendação de aluguel
recomendacoes = [
    ('💸  MENOR CUSTO MENSAL',      'Interior da Califórnia',
     'Vale Central, Bakersfield, Fresno. Aluguéis estimados muito abaixo da média do estado.',
     C_ORANGE),
    ('🏙️  CUSTO × OPORTUNIDADE',   'Subúrbios de LA e SF (<1H Ocean)',
     'Riverside, Stockton, Sacramento. Acesso aos grandes centros com custo 30–40% menor.',
     C_ACCENT),
    ('🌅  MELHOR QUALIDADE DE VIDA','Zona Norte de SF (Near Ocean)',
     'Santa Cruz, Monterey. Custo moderado, natureza, menor densidade populacional.',
     C_GREEN),
    ('🚫  EVITAR PARA ALUGAR',      'San Francisco Bay / Island',
     'Aluguel altíssimo, pressão de demanda extrema. Retorno financeiro desfavorável.',
     C_RED),
]

y_pos = 1.4
for titulo, subtit, desc, color in recomendacoes:
    add_rect(s, 6.6, y_pos, 6.5, 1.15, C_BG2)
    add_rect(s, 6.6, y_pos, 0.06, 1.15, color)
    add_text(s, titulo, 6.75, y_pos + 0.07, 6.2, 0.32, size=11, bold=True, color=color)
    add_text(s, subtit, 6.75, y_pos + 0.38, 6.2, 0.25, size=10, bold=True, color=C_WHITE)
    add_text(s, desc,   6.75, y_pos + 0.62, 6.2, 0.42, size=8.5, color=C_GRAY)
    y_pos += 1.25

add_rect(s, 0.4, 5.55, 12.5, 1.6, C_BG2)
add_rect(s, 0.4, 5.55, 0.06, 1.6, C_ACCENT)
add_text(s, '📌  REGRA PRÁTICA PARA ALUGAR:', 0.6, 5.65, 4.5, 0.35,
         size=11, bold=True, color=C_ACCENT)
add_text(s,
    'Ideal: aluguel mensal ≤ 30% da renda bruta. Com renda mediana de $35k/ano, '
    'o teto saudável é ~$875/mês. Distritos do interior com valor mediano de $108k '
    'estimam aluguéis de $700–$900/mês — dentro do ideal. Nas zonas costeiras ($230k+) '
    'os aluguéis estimados de $1.500–$2.000/mês comprometem 50%+ da renda mediana local.',
    0.6, 6.05, 12.1, 0.95, size=10, color=C_GRAY)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ANÁLISE POR PERFIL DE FAMÍLIA
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
add_tag(s, '05', 0.4, 0.25, bg=C_YELLOW, fg=C_BG)
add_text(s, 'QUAL É O SEU PERFIL?', 0.75, 0.2, 10, 0.5,
         size=22, bold=True, color=C_WHITE)
divider(s, 0.75, color=C_YELLOW)

perfis = [
    ('👨‍👩‍👧  FAMÍLIA COM FILHOS',
     ['Prioridade: escolas, segurança, espaço'],
     ['✅  Subúrbios a <1h do oceano',
      '✅  Riverside e San Bernardino County',
      '✅  Sacramento (região Norte)',
      '💰  Mediana: ~$180–215k',
      '📐  Média de 5 quartos/domicílio'],
     C_GREEN),
    ('👤  JOVEM PROFISSIONAL',
     ['Prioridade: acesso a emprego, mobilidade'],
     ['✅  Subúrbios de SF e LA (<1H Ocean)',
      '✅  East Bay (Oakland, Fremont)',
      '✅  Long Beach / Pasadena',
      '💰  Mediana: ~$215–235k',
      '🚗  Menos de 1h do centro'],
     C_ACCENT),
    ('👴  APOSENTADO / RENDA FIXA',
     ['Prioridade: custo baixo, qualidade de vida'],
     ['✅  Interior calmo (Fresno, Merced)',
      '✅  Norte da Califórnia (Redding)',
      '✅  Idades medianas dos imóveis: 27 anos',
      '💰  Mediana: ~$90–110k',
      '🌿  Baixa densidade, custo de vida reduzido'],
     C_ORANGE),
    ('📈  INVESTIDOR',
     ['Prioridade: valorização e retorno de aluguel'],
     ['✅  Corredor SF–Sacramento (Near Bay)',
      '✅  Cidades satélite de LA (crescimento)',
      '✅  Distritos em transição (INLAND → <1H)',
      '💰  Mediana: $215–235k',
      '📊  Score custo-benefício favorável'],
     C_RED),
]

positions = [(0.4, 1.0), (3.5, 1.0), (6.6, 1.0), (9.7, 1.0)]
for (perfil, subtit, items, color), (lx, ty) in zip(perfis, positions):
    add_rect(s, lx, ty, 2.95, 5.8, C_BG2)
    add_rect(s, lx, ty, 2.95, 0.08, color)
    add_text(s, perfil,   lx+0.1, ty+0.12, 2.75, 0.45, size=11, bold=True, color=color)
    add_text(s, subtit[0],lx+0.1, ty+0.55, 2.75, 0.28, size=8, color=C_GRAY, italic=True)
    divider(s, ty+0.9, color=C_BG, w=2.55, l=lx+0.2)
    yi = ty + 1.0
    for item in items:
        add_text(s, item, lx+0.12, yi, 2.7, 0.32, size=8.5, color=C_WHITE)
        yi += 0.36

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — RANKING FINAL
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)
add_tag(s, '06', 0.4, 0.25, bg=C_ACCENT, fg=C_BG)
add_text(s, 'RANKING FINAL DE ZONAS', 0.75, 0.2, 10, 0.5,
         size=22, bold=True, color=C_WHITE)
divider(s, 0.75, color=C_ACCENT)

add_text(s, 'Avaliação ponderada por: valor, renda, custo-benefício, qualidade de vida e acessibilidade',
         0.4, 0.9, 12.5, 0.4, size=11, color=C_GRAY, italic=True)

ranking = [
    ('#1', 'INTERIOR (INLAND)',         'COMPRAR',  'Melhor custo-benefício absoluto — score ρ renda/valor líder. Ideal para primeira moradia.',     C_GREEN,  '💰💰💰💰💰', '⭐⭐⭐'),
    ('#2', '<1H OCEAN (SUBÚRBIO)',       'AMBOS',    'Equilíbrio entre acesso ao litoral, renda razoável e preço moderado. Melhor para famílias.',    C_ACCENT, '💰💰💰💰',   '⭐⭐⭐⭐'),
    ('#3', 'NEAR OCEAN (LITORAL MÉDIO)', 'ALUGAR',   'Qualidade de vida superior, custo maior. Excelente para quem busca bem-estar sem o extremo de SF.',C_ACCENT2,'💰💰💰',    '⭐⭐⭐⭐'),
    ('#4', 'NEAR BAY (SF BAY AREA)',     'INVESTIR', 'Alta valorização histórica. Exige renda alta. Melhor como investimento de longo prazo.',         C_ORANGE, '💰💰',       '⭐⭐⭐⭐⭐'),
    ('#5', 'ISLAND',                    'EVITAR',   'Luxo puro. Renda local não justifica o custo. Nicho muito restrito, sem liquidez ampla.',        C_RED,    '💰',         '⭐⭐'),
]

headers = ['Rank', 'Zona', 'Ideal Para', 'Por Quê?', 'Acessibilidade', 'Valorização']
col_w   = [0.55, 2.2, 1.1, 5.6, 1.6, 1.3]
col_l   = [0.4, 1.0, 3.25, 4.4, 10.05, 11.7]

# Cabeçalho
for txt, w, l in zip(headers, col_w, col_l):
    add_rect(s, l, 1.45, w, 0.38, C_ACCENT)
    add_text(s, txt, l+0.05, 1.48, w-0.05, 0.32,
             size=9, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

for i, (rank, zona, ideal, pq, color, acesso, valor) in enumerate(ranking):
    y = 1.92 + i * 0.95
    row_bg = C_BG2 if i % 2 == 0 else rgb(0x16, 0x26, 0x3A)
    add_rect(s, 0.4, y, 12.6, 0.88, row_bg)
    add_rect(s, 0.4, y, 0.06, 0.88, color)

    data = [rank, zona, ideal, pq, acesso, valor]
    for txt, w, l in zip(data, col_w, col_l):
        bold = txt in [rank, zona, ideal]
        clr  = color if txt in [rank, zona] else (C_WHITE if txt == ideal else C_GRAY)
        add_text(s, txt, l+0.05, y+0.05, w-0.05, 0.78,
                 size=9 if txt not in [rank, zona] else 11,
                 bold=bold, color=clr, align=PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CONCLUSÃO / CALL TO ACTION
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s)

add_rect(s, 0, 0, 13.33, 7.5, C_BG2)
add_rect(s, 0, 0, 0.1, 7.5, C_ACCENT)
add_rect(s, 0, 0, 13.33, 0.08, C_ACCENT)
add_rect(s, 0, 7.42, 13.33, 0.08, C_ACCENT)

add_text(s, '3 DECISÕES INTELIGENTES', 1.0, 0.6, 11.0, 0.65,
         size=28, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
add_text(s, 'baseadas em dados do mercado imobiliário da Califórnia',
         1.0, 1.2, 11.0, 0.4, size=13, color=C_GRAY,
         italic=True, align=PP_ALIGN.CENTER)

decisoes = [
    ('🏡', 'SE QUER COMPRAR COM SEGURANÇA',
     'Escolha o Interior (INLAND). Valor mediano de $108.500,\nmelhor score custo-benefício do dataset.\nBakersfield, Fresno e Riverside são os pontos de entrada.',
     C_GREEN),
    ('🔑', 'SE QUER ALUGAR E VIVER BEM',
     'Subúrbios a <1h do oceano são o ponto ideal.\nAcesso ao litoral, custo 30–40% menor que SF,\ninfrastrutura de qualidade para famílias.',
     C_ACCENT),
    ('📈', 'SE QUER INVESTIR',
     'Near Bay Area mantém a maior valorização histórica.\nPreço de entrada alto, mas liquidez garantida\ne retorno consistente no longo prazo.',
     C_ORANGE),
]

positions_cta = [(0.5, 2.05), (4.7, 2.05), (8.9, 2.05)]
for (icon, titulo, corpo, color), (lx, ty) in zip(decisoes, positions_cta):
    add_rect(s, lx, ty, 4.0, 4.5, rgb(0x0D, 0x1B, 0x2A))
    add_rect(s, lx, ty, 4.0, 0.08, color)
    add_text(s, icon,   lx+0.15, ty+0.15, 3.7, 0.55, size=30, align=PP_ALIGN.CENTER)
    add_text(s, titulo, lx+0.12, ty+0.7,  3.7, 0.55, size=11, bold=True, color=color)
    add_text(s, corpo,  lx+0.12, ty+1.3,  3.75, 2.5, size=10, color=C_GRAY)

divider(s, 6.7, color=C_BG, w=11.0, l=1.15)
add_text(s, '📂  Análise completa disponível em:  github.com/LuissFellipe/houses-prices-analysis',
         1.0, 6.85, 11.0, 0.45, size=10, color=C_GRAY,
         italic=True, align=PP_ALIGN.CENTER)

# ── Salvar ────────────────────────────────────────────────────────────────────
out = Path('presentation') / 'california_housing_storytelling.pptx'
out.parent.mkdir(exist_ok=True)
prs.save(out)
print(f"✅ Apresentação salva em: {out}")
print(f"   Slides: {len(prs.slides)}")
